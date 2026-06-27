import os
import tempfile
import pandas as pd
import numpy as np

# ReportLab Imports
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak, KeepTogether
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas

# OpenPyXL Imports
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils.dataframe import dataframe_to_rows

# Python-PPTX Imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Local utilities
from utils.visualization import generate_static_plot
from utils.statistics import get_column_types

class NumberedCanvas(canvas.Canvas):
    """
    Two-pass canvas to dynamically add total page counts and footers/headers.
    """
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_decorations(num_pages)
            super().showPage()
        super().save()

    def draw_page_decorations(self, page_count):
        self.saveState()
        
        # Suppress footer on the Cover Page (Page 1)
        if self._pageNumber > 1:
            self.setFont("Helvetica", 9)
            self.setFillColor(colors.HexColor("#475569"))
            
            # Header
            self.drawString(54, 750, "InsightX – Exploratory Data Analysis Report")
            self.setStrokeColor(colors.HexColor("#CBD5E1"))
            self.setLineWidth(0.5)
            self.line(54, 742, 558, 742)
            
            # Footer
            self.drawString(54, 40, "Confidential | InsightX EDA Platform")
            page_text = f"Page {self._pageNumber} of {page_count}"
            self.drawRightString(558, 40, page_text)
            self.line(54, 52, 558, 52)
            
        self.restoreState()


def generate_pdf_report(df, filename, stats, insights, correlations, output_path):
    """
    Generate a highly formatted PDF report using ReportLab.
    """
    # Create the document flow
    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        leftMargin=54,
        rightMargin=54,
        topMargin=72,
        bottomMargin=72
    )
    
    styles = getSampleStyleSheet()
    
    # Custom Palette Styling
    primary_color = colors.HexColor("#4F46E5") # Sleek Indigo
    secondary_color = colors.HexColor("#06B6D4") # Ocean Teal
    text_dark = colors.HexColor("#0F172A") # Slate 900
    text_muted = colors.HexColor("#475569") # Slate 600
    bg_light = colors.HexColor("#F8FAFC") # Slate 50
    
    # Custom styles
    title_style = ParagraphStyle(
        'CoverTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=32,
        leading=38,
        textColor=primary_color,
        alignment=0, # Left
        spaceAfter=12
    )
    
    subtitle_style = ParagraphStyle(
        'CoverSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=18,
        leading=22,
        textColor=secondary_color,
        alignment=0,
        spaceAfter=40
    )
    
    body_style = ParagraphStyle(
        'BodyDark',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10.5,
        leading=14.5,
        textColor=text_dark,
        spaceAfter=10
    )
    
    h1_style = ParagraphStyle(
        'SectionH1',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=22,
        textColor=primary_color,
        spaceBefore=18,
        spaceAfter=10,
        keepWithNext=True
    )
    
    h2_style = ParagraphStyle(
        'SectionH2',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=16,
        textColor=text_dark,
        spaceBefore=12,
        spaceAfter=6,
        keepWithNext=True
    )
    
    meta_style = ParagraphStyle(
        'CoverMetadata',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=text_muted
    )
    
    bullet_style = ParagraphStyle(
        'InsightBullet',
        parent=body_style,
        leftIndent=20,
        firstLineIndent=-10,
        spaceAfter=6
    )

    story = []
    
    # --- PAGE 1: COVER PAGE ---
    story.append(Spacer(1, 100))
    # Brand/Header
    story.append(Paragraph("INSIGHTX", ParagraphStyle('Brand', parent=styles['Normal'], fontName='Helvetica-Bold', fontSize=14, leading=16, textColor=secondary_color, spaceAfter=8)))
    story.append(Paragraph("EXPLORATORY DATA ANALYSIS REPORT", title_style))
    story.append(Paragraph(f"Dataset: {filename}", subtitle_style))
    
    story.append(Spacer(1, 120))
    
    # Metadata Box
    metadata_text = f"""
    <b>Platform:</b> InsightX EDA Platform<br/>
    <b>Date Generated:</b> {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}<br/>
    <b>Record Count:</b> {len(df):,}<br/>
    <b>Feature Count:</b> {len(df.columns)}<br/>
    <b>Generated By:</b> Automated Insight Engine<br/>
    """
    story.append(Paragraph(metadata_text, meta_style))
    story.append(PageBreak())
    
    # --- PAGE 2: EXECUTIVE SUMMARY & METRIC STATS ---
    story.append(Paragraph("Executive Summary", h1_style))
    summary_p = f"""
    This report contains a comprehensive, multi-dimensional analysis of the uploaded dataset <b>{filename}</b>. 
    Using statistics, anomaly profiling, and automated machine learning heuristics, we have audited 
    the quality, distributions, and correlation frameworks of your features to provide deep, actionable insights 
    without manual code authoring.
    """
    story.append(Paragraph(summary_p, body_style))
    
    story.append(Spacer(1, 10))
    story.append(Paragraph("Dataset Overview Metrics", h2_style))
    
    col_types = get_column_types(df)
    
    # Simple table for stats
    overview_data = [
        [Paragraph("<b>Metric</b>", body_style), Paragraph("<b>Value</b>", body_style)],
        [Paragraph("Total Rows", body_style), Paragraph(f"{len(df):,}", body_style)],
        [Paragraph("Total Columns", body_style), Paragraph(f"{len(df.columns)}", body_style)],
        [Paragraph("Numeric Columns", body_style), Paragraph(f"{len(col_types['numeric'])}", body_style)],
        [Paragraph("Categorical Columns", body_style), Paragraph(f"{len(col_types['categorical'])}", body_style)],
        [Paragraph("Boolean Columns", body_style), Paragraph(f"{len(col_types['boolean'])}", body_style)],
        [Paragraph("Date/Time Columns", body_style), Paragraph(f"{len(col_types['datetime'])}", body_style)],
        [Paragraph("Duplicate Rows", body_style), Paragraph(f"{df.duplicated().sum():,}", body_style)],
        [Paragraph("Missing Cells", body_style), Paragraph(f"{df.isna().sum().sum():,}", body_style)]
    ]
    
    t = Table(overview_data, colWidths=[200, 200])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (1,0), primary_color),
        ('TEXTCOLOR', (0,0), (1,0), colors.white),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('BOTTOMPADDING', (0,0), (-1,0), 6),
        ('TOPPADDING', (0,0), (-1,0), 6),
        ('BACKGROUND', (0,1), (-1,-1), bg_light),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#CBD5E1")),
        ('PADDING', (0,0), (-1,-1), 8),
    ]))
    
    # Update textcolor for header cells in table style, but since we used Paragraphs, we must change Paragraph style
    header_style = ParagraphStyle('HeaderStyle', parent=body_style, textColor=colors.white, fontName='Helvetica-Bold')
    overview_data[0] = [Paragraph("Metric", header_style), Paragraph("Value", header_style)]
    
    story.append(t)
    story.append(PageBreak())
    
    # --- PAGE 3: DATA VISUALIZATIONS & PLOTS ---
    story.append(Paragraph("Data Quality & Distribution Charts", h1_style))
    story.append(Paragraph("Visual distributions of dataset missingness map and core features are plotted below:", body_style))
    
    # Generate missingness heatmaps and save locally to embed in report
    try:
        # Create temp files to hold image outputs
        temp_dir = tempfile.gettempdir()
        missing_img_path = os.path.join(temp_dir, 'insightx_missing_map.png')
        dist_img_path = os.path.join(temp_dir, 'insightx_dist_map.png')
        
        # 1. Missingness Map
        missing_bytes = generate_static_plot(df, 'missing_map', {})
        with open(missing_img_path, 'wb') as f:
            f.write(missing_bytes)
        
        img1 = Image(missing_img_path, width=4.5*inch, height=2.8*inch)
        story.append(Paragraph("Missing Values Map Indicator", h2_style))
        story.append(img1)
        story.append(Spacer(1, 15))
        
        # 2. Key Numeric Distribution (First numeric column if available)
        if col_types['numeric']:
            first_num = col_types['numeric'][0]
            dist_bytes = generate_static_plot(df, 'distribution', {'column': first_num})
            with open(dist_img_path, 'wb') as f:
                f.write(dist_bytes)
            
            img2 = Image(dist_img_path, width=4.5*inch, height=2.8*inch)
            story.append(Paragraph(f"Primary Feature Distribution: {first_num}", h2_style))
            story.append(img2)
            
    except Exception as e:
        story.append(Paragraph(f"Error drawing static charts: {str(e)}", body_style))
        
    story.append(PageBreak())
    
    # --- PAGE 4: AUTOMATIC INSIGHTS & FINDINGS ---
    story.append(Paragraph("Automated Insights & Findings", h1_style))
    story.append(Paragraph("Our automated AI analytics engines have parsed your features and extracted these key recommendations and conclusions:", body_style))
    story.append(Spacer(1, 10))
    
    for ins in insights:
        bullet_text = f"• <b>[{ins['category']}]</b>: {ins['text']}"
        story.append(Paragraph(bullet_text, bullet_style))
        
    story.append(Spacer(1, 20))
    story.append(Paragraph("Recommendations & Next Steps", h1_style))
    
    recs = [
        "1. <b>Impute or Drop Null Fields</b>: Use the data cleaning module to handle empty values before training models.",
        "2. <b>Filter Outliers</b>: Standardize feature values to mitigate influence of heavy tail outliers.",
        "3. <b>Encode Categorical Data</b>: Apply one-hot encoding or labeling to categorical variables like top classes for downstream ML pipelines.",
        "4. <b>Mitigate Correlation Redundancy</b>: Address highly correlated feature pairs (above 0.8) by selecting one and discarding the other to reduce multicollinearity."
    ]
    for r in recs:
        story.append(Paragraph(r, body_style))
        
    # Build Document
    doc.build(story, canvasmaker=NumberedCanvas)


def generate_excel_report(df, stats, correlations, output_path):
    """
    Generate a beautifully formatted Multi-Sheet Excel Workbook with openpyxl.
    """
    wb = Workbook()
    
    # Define Styles
    header_fill = PatternFill(start_color="4F46E5", end_color="4F46E5", fill_type="solid")
    stripe_fill = PatternFill(start_color="F8FAFC", end_color="F8FAFC", fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    bold_font = Font(name="Calibri", size=11, bold=True)
    regular_font = Font(name="Calibri", size=11)
    
    thin_border = Border(
        left=Side(style='thin', color='CBD5E1'),
        right=Side(style='thin', color='CBD5E1'),
        top=Side(style='thin', color='CBD5E1'),
        bottom=Side(style='thin', color='CBD5E1')
    )
    
    # --- SHEET 1: Dataset Summary ---
    ws1 = wb.active
    ws1.title = "Overview"
    ws1.views.sheetView[0].showGridLines = True
    
    # Add title block
    ws1.append(["InsightX Exploratory Data Analysis - Overview"])
    ws1.cell(1, 1).font = Font(name="Calibri", size=16, bold=True, color="4F46E5")
    ws1.append([])
    
    overview_headers = ["Metric", "Value"]
    ws1.append(overview_headers)
    for col_idx, text in enumerate(overview_headers, 1):
        cell = ws1.cell(row=3, column=col_idx)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="left")
        
    col_types = get_column_types(df)
    overview_data = [
        ["Dataset Name", "Uploaded Data"],
        ["Total Rows", len(df)],
        ["Total Columns", len(df.columns)],
        ["Numeric Columns", len(col_types['numeric'])],
        ["Categorical Columns", len(col_types['categorical'])],
        ["Boolean Columns", len(col_types['boolean'])],
        ["Date Columns", len(col_types['datetime'])],
        ["Duplicate Rows", int(df.duplicated().sum())],
        ["Missing Cells", int(df.isna().sum().sum())]
    ]
    
    for row_data in overview_data:
        ws1.append(row_data)
        
    # Format overview cells
    for r in range(4, 4 + len(overview_data)):
        for c in range(1, 3):
            cell = ws1.cell(row=r, column=c)
            cell.font = regular_font
            cell.border = thin_border
            if r % 2 == 1:
                cell.fill = stripe_fill
                
    # --- SHEET 2: Statistics Summary ---
    ws2 = wb.create_sheet(title="Descriptive Statistics")
    ws2.views.sheetView[0].showGridLines = True
    
    # Build statistics table
    stats_rows = []
    
    # Compile numeric and categorical headers
    stat_keys = ['type', 'count', 'null_count', 'null_percentage', 'unique_count', 'mean', 'median', 'min', 'max', 'std', 'skewness']
    stat_headers = ["Feature", "Type", "Count", "Null Count", "Null %", "Unique Count", "Mean", "Median", "Min", "Max", "Std Dev", "Skewness"]
    
    ws2.append(["Feature Descriptive Statistics"])
    ws2.cell(1, 1).font = Font(name="Calibri", size=14, bold=True, color="4F46E5")
    ws2.append([])
    
    ws2.append(stat_headers)
    for col_idx, text in enumerate(stat_headers, 1):
        cell = ws2.cell(row=3, column=col_idx)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
        
    row_counter = 4
    for col_name, data in stats.items():
        row_cells = [col_name, data.get('type', 'N/A')]
        
        # Pull each stat key
        row_cells.append(data.get('count', 0))
        row_cells.append(data.get('null_count', 0))
        row_cells.append(round(data.get('null_percentage', 0), 2))
        row_cells.append(data.get('unique_count', 0))
        
        # Numeric statistics
        if data.get('type') == 'numeric':
            row_cells.append(round(data.get('mean', 0), 4))
            row_cells.append(round(data.get('median', 0), 4))
            row_cells.append(round(data.get('min', 0), 4))
            row_cells.append(round(data.get('max', 0), 4))
            row_cells.append(round(data.get('std', 0), 4))
            row_cells.append(round(data.get('skewness', 0), 4))
        else:
            # blanks for non-numeric stats
            row_cells.extend(["N/A"] * 6)
            
        ws2.append(row_cells)
        # Apply style to data rows
        for c in range(1, len(stat_headers) + 1):
            cell = ws2.cell(row=row_counter, column=c)
            cell.font = regular_font
            cell.border = thin_border
            cell.alignment = Alignment(horizontal="center")
            if row_counter % 2 == 1:
                cell.fill = stripe_fill
        row_counter += 1
        
    # --- SHEET 3: Correlation Matrix ---
    ws3 = wb.create_sheet(title="Correlations")
    ws3.views.sheetView[0].showGridLines = True
    
    ws3.append(["Pearson Correlation Matrix"])
    ws3.cell(1, 1).font = Font(name="Calibri", size=14, bold=True, color="4F46E5")
    ws3.append([])
    
    pearson = correlations.get('pearson', {})
    corr_cols = list(pearson.keys())
    
    if corr_cols:
        # Headers: empty cell + column names
        header_row = [""] + corr_cols
        ws3.append(header_row)
        for col_idx, text in enumerate(header_row, 1):
            cell = ws3.cell(row=3, column=col_idx)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
            
        row_counter = 4
        for c1 in corr_cols:
            row_cells = [c1]
            for c2 in corr_cols:
                row_cells.append(round(pearson[c1][c2], 4))
            ws3.append(row_cells)
            
            # Format row
            ws3.cell(row=row_counter, column=1).font = bold_font
            ws3.cell(row=row_counter, column=1).border = thin_border
            ws3.cell(row=row_counter, column=1).alignment = Alignment(horizontal="left")
            
            for c_idx in range(2, len(corr_cols) + 2):
                cell = ws3.cell(row=row_counter, column=c_idx)
                cell.font = regular_font
                cell.border = thin_border
                cell.alignment = Alignment(horizontal="center")
                
                # Apply simple color scales
                val = cell.value
                if val is not None:
                    if val > 0.7: # high positive
                        cell.fill = PatternFill(start_color="C6F6D5", end_color="C6F6D5", fill_type="solid") # green
                    elif val < -0.7: # high negative
                        cell.fill = PatternFill(start_color="FED7D7", end_color="FED7D7", fill_type="solid") # red
            row_counter += 1
            
    # Auto-adjust column widths
    for ws in [ws1, ws2, ws3]:
        for col in ws.columns:
            max_len = 0
            for cell in col:
                val_str = str(cell.value or '')
                if cell.row == 1: # Ignore title length for column width calculation
                    continue
                if len(val_str) > max_len:
                    max_len = len(val_str)
            col_letter = col[0].column_letter
            ws.column_dimensions[col_letter].width = max(max_len + 3, 12)
            
    wb.save(output_path)


def generate_pptx_report(df, filename, stats, insights, output_path):
    """
    Generate a modern, clean PowerPoint presentation file.
    """
    prs = Presentation()
    
    # Page setup (widescreen 16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Palette Colors
    c_indigo = "4F46E5"
    c_teal = "06B6D4"
    c_slate = "0F172A"
    c_white = "FFFFFF"
    
    # Helper: Set slide background color
    def set_bg_color(slide, hex_color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = colors.HexColor(f"#{hex_color}").rgb
        
    # Helper: Create title and textbox
    def add_slide_title(slide, text, color=c_indigo):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = colors.HexColor(f"#{color}").rgb
        
    # --- SLIDE 1: Title Slide ---
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    set_bg_color(slide1, c_slate)
    
    # Add large title
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.83), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "INSIGHTX"
    p.font.name = "Arial"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = colors.HexColor(f"#{c_teal}").rgb
    
    p2 = tf.add_paragraph()
    p2.text = "Exploratory Data Analysis & Insights Report"
    p2.font.name = "Arial"
    p2.font.size = Pt(28)
    p2.font.color.rgb = colors.HexColor(f"#{c_white}").rgb
    
    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.83), Inches(1.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = f"Dataset: {filename}\nGenerated automatically on {pd.Timestamp.now().strftime('%Y-%m-%d')}"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = colors.HexColor("#94A3B8").rgb
    
    # --- SLIDE 2: Dataset Overview ---
    slide2 = prs.slides.add_slide(slide_layout)
    set_bg_color(slide2, c_white)
    add_slide_title(slide2, "Dataset Overview Metrics")
    
    # Create Table on slide
    rows, cols = 8, 2
    left = Inches(1.5)
    top = Inches(1.6)
    width = Inches(10.33)
    height = Inches(4.5)
    
    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(6.0)
    table.columns[1].width = Inches(4.33)
    
    col_types = get_column_types(df)
    metrics = [
        ["Overview Metric", "Value"],
        ["Total Rows", f"{len(df):,}"],
        ["Total Columns", f"{len(df.columns)}"],
        ["Numeric Features", f"{len(col_types['numeric'])}"],
        ["Categorical Features", f"{len(col_types['categorical'])}"],
        ["Boolean Features", f"{len(col_types['boolean'])}"],
        ["Null/Missing Values", f"{df.isna().sum().sum():,}"],
        ["Duplicates Detected", f"{df.duplicated().sum():,}"]
    ]
    
    for r_idx, row in enumerate(metrics):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            # Formatting
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT
            
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = colors.HexColor(f"#{c_white}").rgb
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors.HexColor(f"#{c_indigo}").rgb
            else:
                p.font.color.rgb = colors.HexColor(f"#{c_slate}").rgb
                if r_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = colors.HexColor("#F8FAFC").rgb
                    
    # --- SLIDE 3: Key Insights ---
    slide3 = prs.slides.add_slide(slide_layout)
    set_bg_color(slide3, c_white)
    add_slide_title(slide3, "Automated Business Insights")
    
    # Textbox for insights
    insight_box = slide3.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.33), Inches(4.8))
    tf_ins = insight_box.text_frame
    tf_ins.word_wrap = True
    
    for i_idx, ins in enumerate(insights[:6]): # Limit to top 6 insights
        p = tf_ins.add_paragraph() if i_idx > 0 else tf_ins.paragraphs[0]
        # Remove raw HTML tags from text
        clean_text = ins['text'].replace("<b>", "").replace("</b>", "")
        p.text = f"• [{ins['category']}] {clean_text}"
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.space_after = Pt(14)
        p.font.color.rgb = colors.HexColor(f"#{c_slate}").rgb
        
    # --- SLIDE 4: Key Recommendations ---
    slide4 = prs.slides.add_slide(slide_layout)
    set_bg_color(slide4, c_white)
    add_slide_title(slide4, "Recommended Actions")
    
    recs_box = slide4.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(4.5))
    tf_recs = recs_box.text_frame
    tf_recs.word_wrap = True
    
    recs = [
        ("Handle Missing Fields", "Impute null columns or remove them to maintain dataset validity."),
        ("Outlier Treatment", "Apply clipping (winsorization) or drop columns with high anomaly rates."),
        ("Multicollinearity Check", "Select single dominant features from pairs with high (>0.85) correlation to reduce noise."),
        ("Type Conversion Auditing", "Fix schema errors by casting dates and categories correctly on the Analysis tab.")
    ]
    
    for idx, (title, desc) in enumerate(recs):
        p = tf_recs.add_paragraph() if idx > 0 else tf_recs.paragraphs[0]
        p.text = f"{idx+1}. {title}: {desc}"
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True if idx == 0 else False
        p.space_after = Pt(20)
        p.font.color.rgb = colors.HexColor(f"#{c_slate}").rgb
        # highlight heading
        run = p.runs[0]
        run.font.bold = True
        run.font.color.rgb = colors.HexColor(f"#{c_indigo}").rgb
        
    prs.save(output_path)
